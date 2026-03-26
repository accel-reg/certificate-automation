# Certificate Generator GUI
# Author: Accel

import os
import re
import pandas as pd
import pdfplumber
import win32com.client
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import tkinter as tk
from tkinter import ttk, filedialog, messagebox, scrolledtext
import threading


class CertificateGeneratorGUI:
    def __init__(self, root):
        self.root = root
        self.root.title("Certificate Generator")
        self.root.geometry("700x650")
        self.root.resizable(False, False)

        self.template_path = tk.StringVar()
        self.names_file_path = tk.StringVar()
        self.output_folder_path = tk.StringVar()

        self.font_name = tk.StringVar(value="Lavonia Classy")
        self.placeholder_text = tk.StringVar(value="names")
        self.max_font_size = tk.IntVar(value=60)
        self.min_font_size = tk.IntVar(value=36)

        self.names_list = []

        self.setup_ui()

    # ---------------- UI ---------------- #

    def setup_ui(self):
        main_frame = ttk.Frame(self.root, padding="20")
        main_frame.grid(row=0, column=0)

        ttk.Label(main_frame,
                  text="🎓 Certificate Generator",
                  font=('Helvetica', 16, 'bold')
                  ).grid(row=0, column=0, columnspan=3, pady=20)

        file_frame = ttk.LabelFrame(main_frame, text="File Selection", padding=10)
        file_frame.grid(row=1, column=0, columnspan=3, pady=10)

        ttk.Label(file_frame, text="Template PPT:").grid(row=0, column=0)
        ttk.Entry(file_frame, textvariable=self.template_path, width=50).grid(row=0, column=1)
        ttk.Button(file_frame, text="Browse", command=self.browse_template).grid(row=0, column=2)

        ttk.Label(file_frame, text="Names File:").grid(row=1, column=0)
        ttk.Entry(file_frame, textvariable=self.names_file_path, width=50).grid(row=1, column=1)
        ttk.Button(file_frame, text="Browse", command=self.browse_names_file).grid(row=1, column=2)

        ttk.Label(file_frame, text="Output Folder:").grid(row=2, column=0)
        ttk.Entry(file_frame, textvariable=self.output_folder_path, width=50).grid(row=2, column=1)
        ttk.Button(file_frame, text="Browse", command=self.browse_output_folder).grid(row=2, column=2)

        settings_frame = ttk.LabelFrame(main_frame, text="Settings", padding=10)
        settings_frame.grid(row=2, column=0, columnspan=3, pady=10)

        ttk.Label(settings_frame, text="Placeholder Text:").grid(row=0, column=0)
        ttk.Entry(settings_frame, textvariable=self.placeholder_text, width=20).grid(row=0, column=1)

        ttk.Label(settings_frame, text="Font Name:").grid(row=1, column=0)
        ttk.Entry(settings_frame, textvariable=self.font_name, width=20).grid(row=1, column=1)

        ttk.Label(settings_frame, text="Max Font Size:").grid(row=0, column=2)
        ttk.Spinbox(settings_frame, from_=20, to=100,
                    textvariable=self.max_font_size, width=5).grid(row=0, column=3)

        ttk.Label(settings_frame, text="Min Font Size:").grid(row=1, column=2)
        ttk.Spinbox(settings_frame, from_=10, to=60,
                    textvariable=self.min_font_size, width=5).grid(row=1, column=3)

        button_frame = ttk.Frame(main_frame)
        button_frame.grid(row=3, column=0, columnspan=3, pady=10)

        self.load_btn = ttk.Button(button_frame, text="Load Names",
                                   command=self.load_names)
        self.load_btn.grid(row=0, column=0, padx=5)

        self.generate_btn = ttk.Button(button_frame,
                                       text="Generate PDF Certificates",
                                       command=self.start_generation,
                                       state=tk.DISABLED)
        self.generate_btn.grid(row=0, column=1, padx=5)

        ttk.Button(button_frame, text="Exit",
                   command=self.root.quit).grid(row=0, column=2)

        self.count_label = ttk.Label(main_frame,
                                     text="Names Loaded: 0",
                                     font=('Helvetica', 10, 'bold'))
        self.count_label.grid(row=4, column=0, columnspan=3)

        log_frame = ttk.LabelFrame(main_frame, text="Progress Log")
        log_frame.grid(row=5, column=0, columnspan=3)

        self.log_text = scrolledtext.ScrolledText(log_frame,
                                                  width=80, height=12)
        self.log_text.grid(row=0, column=0)

        self.progress = ttk.Progressbar(main_frame,
                                        mode='determinate',
                                        length=660)
        self.progress.grid(row=6, column=0, columnspan=3, pady=10)

        ttk.Label(main_frame,
                  text="© Copyright by Accel",
                  font=('Helvetica', 9, 'italic'),
                  foreground="gray").grid(row=7, column=0, columnspan=3)

    # ---------------- Browse ---------------- #

    def browse_template(self):
        path = filedialog.askopenfilename(filetypes=[("PPTX", "*.pptx")])
        if path:
            self.template_path.set(path)

    def browse_names_file(self):
        path = filedialog.askopenfilename(
            filetypes=[("Excel", "*.xlsx *.xls"), ("PDF", "*.pdf")])
        if path:
            self.names_file_path.set(path)

    def browse_output_folder(self):
        path = filedialog.askdirectory()
        if path:
            self.output_folder_path.set(path)

    # ---------------- Logging ---------------- #

    def log(self, msg):
        self.log_text.insert(tk.END, msg + "\n")
        self.log_text.see(tk.END)
        self.root.update_idletasks()

    # ---------------- Name Extraction ---------------- #

    def extract_names(self, file_path):
        ext = os.path.splitext(file_path)[1].lower()
        names = []

        if ext in [".xlsx", ".xls"]:
            df = pd.read_excel(file_path)

            name_column = None
            for col in df.columns:
                if "name" in str(col).lower():
                    name_column = col
                    break

            if name_column:
                names = df[name_column].dropna().astype(str).tolist()
            else:
                names = df.iloc[:, 0].dropna().astype(str).tolist()

        elif ext == ".pdf":
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    text = page.extract_text()
                    if text:
                        names.extend(text.split("\n"))

        cleaned = []
        for name in names:
            name = re.sub(r"[^a-zA-Z\s.]", "", name)
            name = re.sub(r"\s+", " ", name).strip()
            if len(name) > 2:
                cleaned.append(name)

        return list(dict.fromkeys(cleaned))

    def load_names(self):
        if not self.names_file_path.get():
            messagebox.showerror("Error", "Select names file first.")
            return

        self.names_list = self.extract_names(self.names_file_path.get())

        if not self.names_list:
            messagebox.showwarning("Warning", "No valid names found.")
            return

        self.count_label.config(text=f"Names Loaded: {len(self.names_list)}")
        self.generate_btn.config(state=tk.NORMAL)
        self.log(f"Loaded {len(self.names_list)} names successfully.")

    # ---------------- Font Logic ---------------- #

    def get_font_size(self, name):
        length = len(name)
        if length <= 12:
            return self.max_font_size.get()
        elif length <= 20:
            return int(self.max_font_size.get() * 0.8)
        return self.min_font_size.get()

    # ---------------- PPT → PDF ---------------- #

    def convert_ppt_to_pdf(self, ppt_path, pdf_path, powerpoint):
        presentation = None
        try:
            presentation = powerpoint.Presentations.Open(
                os.path.abspath(ppt_path),
                WithWindow=False
            )
            presentation.SaveAs(os.path.abspath(pdf_path), 32)
        finally:
            if presentation:
                presentation.Close()

    # ---------------- Generate ---------------- #

    def generate_certificates(self):
        try:
            template = self.template_path.get()
            output = self.output_folder_path.get()

            if not template or not output:
                messagebox.showerror("Error", "Select template and output folder.")
                return

            os.makedirs(output, exist_ok=True)

            self.generate_btn.config(state=tk.DISABLED)
            self.load_btn.config(state=tk.DISABLED)

            powerpoint = win32com.client.DispatchEx("PowerPoint.Application")

            self.progress["maximum"] = len(self.names_list)

            for i, name in enumerate(self.names_list, 1):

                prs = Presentation(template)
                slide = prs.slides[0]

                for shape in slide.shapes:
                    if shape.has_text_frame and \
                            self.placeholder_text.get().lower() in shape.text.lower():

                        shape.text_frame.clear()
                        p = shape.text_frame.paragraphs[0]
                        p.text = name
                        p.alignment = PP_ALIGN.CENTER

                        run = p.runs[0]
                        run.font.name = self.font_name.get()
                        run.font.size = Pt(self.get_font_size(name))
                        run.font.bold = True
                        run.font.color.rgb = RGBColor(0, 0, 0)
                        break

                safe = re.sub(r'[<>:"/\\|?*]', '', name).replace(" ", "_")
                ppt_path = os.path.join(output, safe + ".pptx")
                pdf_path = os.path.join(output, safe + ".pdf")

                prs.save(ppt_path)
                prs = None

                self.convert_ppt_to_pdf(ppt_path, pdf_path, powerpoint)

                if os.path.exists(pdf_path):
                    os.remove(ppt_path)

                self.progress["value"] = i
                self.log(f"[{i}/{len(self.names_list)}] Created {safe}.pdf")

            powerpoint.Quit()

            messagebox.showinfo("Success", "All certificates generated.")

        except Exception as e:
            messagebox.showerror("Error", str(e))

        finally:
            self.generate_btn.config(state=tk.NORMAL)
            self.load_btn.config(state=tk.NORMAL)

    def start_generation(self):
        threading.Thread(target=self.generate_certificates,
                         daemon=True).start()


if __name__ == "__main__":
    root = tk.Tk()
    app = CertificateGeneratorGUI(root)
    root.mainloop()